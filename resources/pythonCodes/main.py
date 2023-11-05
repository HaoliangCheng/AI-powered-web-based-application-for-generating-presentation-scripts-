# updated by Difei 07/02/23
# main entry .py file

import argparse
import os
import numpy as np
from pptx import Presentation
import cv2
from pptx.enum.shapes import MSO_SHAPE_TYPE

# TODO: integrate image2txt.py and reorganize the data into json/dictionary

def iter_picture_slide(slide):
    # NOTE: picture is a type of shape
    pictures = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append(shape)
    return pictures

def extract_images_text(prs, output_dir):
    """

    :param prs:
    :param output_dir:
    :return:
    """

    for i in range(len(prs.slides)):
        slide = prs.slides[i]

        # output stream for text
        list_of_elements = []
        text_fname = open(os.path.join(output_dir, f"Text{i+1}.txt"), "w")

        # extract images from slide
        pictures = iter_picture_slide(slide)
        if len(pictures) > 1:
            for j in range(len(pictures)):
                picture = pictures[j]
                image = picture.image
                image_bytes = image.blob
                image_ext = image.ext
                output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
                image_fname = f'Slide{i+1}_{j}.{image_ext}'
                cv2.imwrite(f'./{output_dir}/{image_fname}', output_image)
        elif len(pictures) == 1:
            image = pictures[0].image
            image_bytes = image.blob
            image_ext = image.ext
            output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
            image_fname = f'Slide{i+1}.{image_ext}'
            cv2.imwrite(f'./{output_dir}/{image_fname}', output_image)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                # other types of shape (e.g., hyperlink, video...)
                continue
            else:
                # extract the text
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    line = ''
                    for run in paragraph.runs:
                        if len(line):
                            line += ' ' + run.text
                        else:
                            line += run.text
                    list_of_elements.append(line)

        for elements in list_of_elements:
            text_fname.write(elements + '\n')
        text_fname.close()

def main():
    parser = argparse.ArgumentParser(description='Extract all images and text from a PowerPoint presentation')

    # temp solution
    parser.add_argument("input")
    parser.add_argument("output")

    # # required arguments
    # required_args = parser.add_argument_group('required input arguments')
    # required_args.add_argument('-i', '--input', type=str, required=True,
    #                            help='Path of the input pptx file')
    # # absolute directory path for output
    # required_args.add_argument('-o', '--output', type=str, required=False, default=os.getcwd(),
    #                            help='Name of the directory for extraction output')

    # TODO: Link the variables to Java Backend 
    args = parser.parse_args()
    prs = Presentation(args.input)
    extract_images_text(prs, args.output)

if __name__ == '__main__':
    main()
