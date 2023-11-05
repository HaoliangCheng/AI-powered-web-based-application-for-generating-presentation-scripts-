# updated by Difei 07/02/23
# main entry .py file

import argparse
import os
import cv2
import numpy as np
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from transformers import pipeline

def image2text(img):
    """
    This function extracts the text from input image, which now supports OCR and image caption
    :param img: image path to be parsed
    :return: extracted text from image
    """

    from transformers import pipeline
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel

    # Image caption: a description of the image
    captioner = pipeline("image-to-text", model="Salesforce/blip-image-captioning-base")
    img_caption_txt = captioner(img)[0]['generated_text']

    # TODO: OCR (Optical Character Recognition)

    return img_caption_txt


def iter_picture_slide(slide):
    # NOTE: picture is a type of shape
    pictures = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append(shape)
    return pictures


def extract_images_text_each_slide(prs, output_dir):
    """This function is a deprecated version of extracting content for each slide

    :param prs:
    :param output_dir:
    :return:
    """

    # TODO: integrate image2txt.py and reorganize the data into json/dictionary
    for i in range(len(prs.slides)):
        contents = dict()
        contents['image'] = []

        slide = prs.slides[i]

        # output stream for text
        list_of_elements = []

        # TODO: smart way of storing results
        content_fname = os.path.join(output_dir, f"Content{i + 1}.json")

        ###text_fname = open(os.path.join(output_dir, f"Text{i+1}.txt"), "w")

        # extract images from slide
        pictures = iter_picture_slide(slide)
        if len(pictures) > 1:
            for j in range(len(pictures)):
                picture = pictures[j]
                image = picture.image
                image_bytes = image.blob
                image_ext = image.ext
                output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
                image_fname = f'Slide{i + 1}_{j}.{image_ext}'  # basename
                image_pname = f'{output_dir}/{image_fname}'  # pathname

                contents['image'].append(image_pname)
                cv2.imwrite(image_pname, output_image)
        elif len(pictures) == 1:
            image = pictures[0].image
            image_bytes = image.blob
            image_ext = image.ext
            output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
            image_fname = f'Slide{i + 1}.{image_ext}'
            image_pname = f'{output_dir}/{image_fname}'

            contents['image'].append(image_pname)
            cv2.imwrite(image_pname, output_image)

        # extract text from image if exists
        images = contents['image']
        if len(images) > 0:
            for img in images:
                # TODO: The image is ...
                txt = image2text(img)
                contents[f'image{images.index(img) + 1}'] = txt
        del contents['image']

        # extract text from slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                # other types of shape (e.g., hyperlink, video...)
                continue
            else:
                # extract the text
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    line = ''
                    for run in paragraph.runs:
                        if len(line):
                            line += ' ' + run.text
                        else:
                            line += run.text
                    list_of_elements.append(line)

        contents['text'] = ''
        for elements in list_of_elements:
            contents['text'] += elements + '\n'

        contents_json = json.dumps(contents, indent=4)
        with open(content_fname, "w") as outfile:
            outfile.write(contents_json)
        outfile.close()


def extract_images_text(prs, output_dir):
    """This function is designed to extract the content from input slides

    :param prs:
    :param output_dir:
    :return:
    """

    contents = dict()
    for i in range(len(prs.slides)):
        content = dict()
        content['image'] = []

        slide = prs.slides[i]

        # output stream for text
        list_of_elements = []

        ###text_fname = open(os.path.join(output_dir, f"Text{i+1}.txt"), "w")

        # extract images from slide
        pictures = iter_picture_slide(slide)
        if len(pictures) > 1:
            for j in range(len(pictures)):
                picture = pictures[j]
                image = picture.image
                image_bytes = image.blob
                image_ext = image.ext
                output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
                image_fname = f'Slide{i + 1}_{j}.{image_ext}'  # basename
                image_pname = f'{output_dir}/{image_fname}'  # pathname

                content['image'].append(image_pname)
                cv2.imwrite(image_pname, output_image)
        elif len(pictures) == 1:
            image = pictures[0].image
            image_bytes = image.blob
            image_ext = image.ext
            output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
            image_fname = f'Slide{i + 1}.{image_ext}'
            image_pname = f'{output_dir}/{image_fname}'

            content['image'].append(image_pname)
            cv2.imwrite(image_pname, output_image)

        # extract text from image if exists
        images = content['image']
        if len(images) > 0:
            for img in images:
                # TODO: The image is ...
                txt = image2text(img)
                content[f'image{images.index(img) + 1}'] = txt
        del content['image']

        # extract text from slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                # other types of shape (e.g., hyperlink, video...)
                continue
            else:
                # extract the text
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    line = ''
                    for run in paragraph.runs:
                        if len(line):
                            line += ' ' + run.text
                        else:
                            line += run.text
                    list_of_elements.append(line)

        content['text'] = ''
        for elements in list_of_elements:
            content['text'] += elements + '\n'

        contents[f"Slide{i}"] = content

    content_fname = os.path.join(output_dir, "Content.json")
    contents_json = json.dumps(contents, indent=4)
    with open(content_fname, "w") as outfile:
        outfile.write(contents_json)
    outfile.close()


def main():
    parser = argparse.ArgumentParser(description='Extract all images and text from a PowerPoint presentation')
    parser.add_argument("input")
    parser.add_argument("output")

    # required arguments
    # required_args = parser.add_argument_group('required input arguments')
    # required_args.add_argument('-i', '--input', type=str, required=True,
    #                            help='Path of the input pptx file')
    # absolute directory path for output
    # required_args.add_argument('-o', '--output', type=str, required=True,
    #                            help='Name of the directory for extraction output')

    args = parser.parse_args()
    prs = Presentation(args.input)

    try:
        extract_images_text(prs, args.output)
    except:
        raise ValueError('Fail to parse the input .pptx file.')


if __name__ == '__main__':
    main()